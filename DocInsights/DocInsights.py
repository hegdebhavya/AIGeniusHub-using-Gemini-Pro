import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

def get_text_from_file(file):
    text = ""
    if file.type == "application/pdf":
        pdf_reader = PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        ppt = Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, you should also be able to summarize the entire content in the file. if the answer is not in
    provided context just say, 'answer is not available in the context', don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)

    chain = get_conversational_chain()
    
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    output_text = response["output_text"]  # Accessing the output_text key from the response
    print(output_text)
    st.write("Reply: ", output_text)  # Displaying only the text part of the response


def main():
    st.set_page_config("Chat with Files")
    st.header('Extract Answers from Documents:mag_right:', divider='rainbow')
    user_question = st.text_input("Ask a Question from Document")
    ask_button = st.button("Ask")
    if ask_button and user_question:
        user_input(user_question)
    with st.sidebar:
        st.title("Embedding creation")
        uploaded_files = st.file_uploader("Upload your Files and Click on the Submit & Process Button", accept_multiple_files=True, type=['pdf', 'docx', 'pptx'])
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                raw_text = ""
                for file in uploaded_files:
                    raw_text += get_text_from_file(file)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("Done")

if __name__ == "__main__":
    main()
